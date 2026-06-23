# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r"e:\文档\选题\答辩PPT\开题答辩-何至恒.pptx"
prs = Presentation(path)

def walk(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(sh.shapes)

with open(r"e:\文档\选题\答辩PPT\_verify.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides, 1):
        f.write(f"\n=== PAGE {i} ===\n")
        for sh in walk(slide.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                f.write(sh.text_frame.text.strip()[:200] + "\n---\n")
        pics = sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
        f.write(f"[pictures: {pics}]\n")
print("ok")
