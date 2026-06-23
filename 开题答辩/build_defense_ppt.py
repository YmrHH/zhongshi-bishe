# -*- coding: utf-8 -*-
"""按《开题答辩PPT填写大纲.md》14 页标准版，从模板生成答辩 PPT。

修改本文件中的幻灯片文案后，请同步更新同目录下《开题答辩PPT填写大纲.md》§三、§四；
若影响讲解内容，请同步核对《开题答辩指南.md》§四—§六。
"""

from __future__ import annotations

import os
import shutil
from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
ARCH_FILL = RGBColor(0xD6, 0xE8, 0xF7)
ARCH_LINE = RGBColor(0x2E, 0x75, 0xB6)
CARD_FILL = RGBColor(0xF5, 0xF7, 0xFA)
CARD_LINE = RGBColor(0x2E, 0x75, 0xB6)

DIR = os.path.dirname(__file__)
ROOT = os.path.dirname(DIR)
TEMPLATE = os.path.join(DIR, "开题报告答辩模板.pptx")
OUTPUT = os.path.join(DIR, "开题答辩-何至恒.pptx")

STUDENT = "何至恒"
STUDENT_ID = "1230120215"
ADVISOR = "吕志峰"
DATE = "2026年6月"

ACTIVITY_IMG = os.path.join(ROOT, "业务活动图", "png", "1中试需求管理业务活动图.png")
MODULE_IMG = os.path.join(
    ROOT, "功能模块图", "广州生产力促进中心中试服务管理系统的设计与实现.png"
)

TITLE_HEADERS = {"论文选题背景", "论文表达结构", "论文研究方法", "论文主要结论"}
PLACEHOLDER_MARKERS = (
    "单击此处",
    "在此输入",
    "点此编辑",
    "添加标题",
    "内容要复合",
    "内容要与标题",
    "PPT下载",
    "1ppt.com",
)

# 保留的模板页（0-based）；27 为复制的第 23 页
KEEP_ORIG = [0, 1, 7, 5, 15, 17, 28, 10, 22, 11, 27, 13, 23, 24]
TARGET_ORDER = [0, 1, 7, 5, 15, 17, 28, 10, 22, 11, 27, 13, 23, 24]


# ── 工具函数 ──────────────────────────────────────────────


def walk_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(sh.shapes)


def iter_text_shapes(shapes, off_left=0, off_top=0):
    """遍历文本 shape，返回 (shape, abs_left, abs_top)。"""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_shapes(sh.shapes, off_left + sh.left, off_top + sh.top)
        elif sh.has_text_frame:
            yield sh, off_left + sh.left, off_top + sh.top


def _is_title_text(t: str) -> bool:
    return t in ("添加标题", "点此编辑标题")


def _is_body_text(t: str) -> bool:
    return (
        "单击" in t or "内容要复合" in t or "内容要与标题" in t or "在此输入" in t
    ) and not _is_title_text(t)


def collect_card_boxes(slide, min_abs_left=0):
    """配对标题/正文：优先顶层文本框；否则按 GROUP；最后扁平遍历。"""
    top_titles, top_bodies = [], []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP or not sh.has_text_frame:
            continue
        al, at = sh.left, sh.top
        if al < min_abs_left:
            continue
        t = sh.text_frame.text.strip()
        if _is_title_text(t):
            top_titles.append((at, al, sh))
        elif _is_body_text(t):
            top_bodies.append((at, al, sh))
    if len(top_titles) >= 2 and len(top_titles) == len(top_bodies):
        top_titles.sort(key=lambda x: (x[0] // 400000, x[1]))
        top_bodies.sort(key=lambda x: (x[0] // 400000, x[1]))
        return [s for _, _, s in top_titles], [s for _, _, s in top_bodies]

    groups = []
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        titles, bodies = [], []
        for sub, al, at in iter_text_shapes([sh], sh.left, sh.top):
            if al < min_abs_left:
                continue
            t = sub.text_frame.text.strip()
            if _is_title_text(t):
                titles.append(sub)
            elif _is_body_text(t):
                bodies.append(sub)
        if titles and bodies:
            groups.append((sh.top, sh.left, titles[0], bodies[0]))

    if groups:
        groups.sort(key=lambda g: (g[0] // 500000, g[1]))
        return [g[2] for g in groups], [g[3] for g in groups]

    seen = set()
    titles, bodies = [], []
    for sh, abs_left, abs_top in iter_text_shapes(slide.shapes):
        if abs_left < min_abs_left:
            continue
        key = id(sh._element)
        if key in seen:
            continue
        seen.add(key)
        t = sh.text_frame.text.strip()
        if _is_title_text(t):
            titles.append((abs_top, abs_left, sh))
        elif _is_body_text(t):
            bodies.append((abs_top, abs_left, sh))
    titles.sort(key=lambda x: (x[0] // 500000, x[1]))
    bodies.sort(key=lambda x: (x[0] // 500000, x[1]))
    return [s for _, _, s in titles], [s for _, _, s in bodies]


def scrub_shape_text(slide, substring: str, replacement: str = ""):
    for sh in walk_shapes(slide.shapes):
        if sh.has_text_frame and substring in sh.text_frame.text:
            set_shape_text(sh, replacement, 12)


def scrub_placeholders(slide):
    """清除仍含模板占位符文字的 shape。"""
    for sh in walk_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if any(m in t for m in PLACEHOLDER_MARKERS):
            set_shape_text(sh, "", 12)


def set_shape_text(shape, text, size=22, bold=False, color=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = "微软雅黑"
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_section_bullet_list(shape, sections, header_size=18, bullet_size=17):
    """条目列表：[(小标题, [要点…]), …]。"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    first = True
    for header, bullets in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = header
        r.font.size = Pt(header_size)
        r.font.name = "微软雅黑"
        r.font.bold = True
        for item in bullets:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = f"• {item}"
            r.font.size = Pt(bullet_size)
            r.font.name = "微软雅黑"


def set_bullet_list(shape, items, size=16, align=PP_ALIGN.LEFT, bullet="•"):
    """单栏条目列表。"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = f"{bullet} {item}"
        r.font.size = Pt(size)
        r.font.name = "微软雅黑"


def set_title_bullet_box(
    shape, title, bullets, title_size=20, bullet_size=15, body_align=PP_ALIGN.LEFT
):
    """标题居中 + 正文条目列表（圆角卡片等）。"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(title_size)
    r1.font.name = "微软雅黑"
    r1.font.bold = True
    for item in bullets:
        p = tf.add_paragraph()
        p.alignment = body_align
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(bullet_size)
        r.font.name = "微软雅黑"


def fill_shape_body(shape, body, size=16):
    """正文：字符串或条目列表。"""
    if isinstance(body, (list, tuple)):
        set_bullet_list(shape, body, size=size)
    else:
        set_shape_text(shape, body, size=size)


def set_card_compact(shape, title: str, body: str, title_size=14, body_size=12):
    """窄卡片：标题 + 短正文（同形状两行，避免模板双 FREEFORM 重叠）。"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(4)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(title_size)
    r1.font.name = "微软雅黑"
    r1.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(body_size)
    r2.font.name = "微软雅黑"


def replace_exact(slide, old: str, new: str, size=22, bold=False):
    for sh in walk_shapes(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip() == old:
            set_shape_text(sh, new, size=size, bold=bold)


def replace_contains(slide, old: str, new: str, size=22, bold=False):
    for sh in walk_shapes(slide.shapes):
        if sh.has_text_frame and old in sh.text_frame.text:
            set_shape_text(sh, new, size=size, bold=bold)


def set_page_title(slide, title: str):
    for sh in walk_shapes(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip() in TITLE_HEADERS:
            set_shape_text(sh, title, size=34, bold=True)
            return
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top < 600000:
            set_shape_text(sh, title, size=34, bold=True)
            return


def shapes_with_text(slide, substring: str):
    return [
        sh
        for sh in walk_shapes(slide.shapes)
        if sh.has_text_frame and substring in sh.text_frame.text
    ]


def remove_top_level_containing(slide, substring: str):
    for sh in list(slide.shapes):
        for sub in walk_shapes([sh]):
            if sub.has_text_frame and substring in sub.text_frame.text:
                sh.element.getparent().remove(sh.element)
                break


def fill_by_marker(slide, marker: str, text: str, size=20):
    for sh in walk_shapes(slide.shapes):
        if sh.has_text_frame and marker in sh.text_frame.text:
            set_shape_text(sh, text, size=size)


def fill_brick_five_cards(slide, cards, row_split=Inches(8.0), body_size=13):
    """模板第 11 页：上排 2 卡 + 下排 3 卡，按业务顺序填标题与条目正文。"""
    titles, bodies = [], []
    for sh, al, at in iter_text_shapes(slide.shapes):
        t = sh.text_frame.text.strip()
        if _is_title_text(t):
            titles.append((at, al, sh))
        elif _is_body_text(t):
            bodies.append((at, al, sh))
    if len(titles) < 5 or len(bodies) < 5:
        fill_card_pairs(slide, cards, body_size=body_size)
        return

    def row_sort(boxes):
        top_row = sorted([b for b in boxes if b[0] < row_split], key=lambda x: x[1])
        bottom_row = sorted([b for b in boxes if b[0] >= row_split], key=lambda x: x[1])
        return [s for _, _, s in top_row] + [s for _, _, s in bottom_row]

    title_boxes = row_sort(titles)
    body_boxes = row_sort(bodies)
    for i, (title, body) in enumerate(cards):
        if i < len(title_boxes):
            set_shape_text(title_boxes[i], title, size=18, bold=True)
        if i < len(body_boxes):
            fill_shape_body(body_boxes[i], body, size=body_size)
    for sh in title_boxes[len(cards) :]:
        set_shape_text(sh, "", 14)
    for sh in body_boxes[len(cards) :]:
        set_shape_text(sh, "", 14)


def fill_card_pairs(slide, pairs, min_abs_left=0, body_size=18):
    """按绝对位置填充 (title, body) 卡片对。"""
    title_boxes, body_boxes = collect_card_boxes(slide, min_abs_left)
    for i, (title, body) in enumerate(pairs):
        if i < len(title_boxes):
            set_shape_text(title_boxes[i], title, size=20, bold=True)
        if i < len(body_boxes):
            fill_shape_body(body_boxes[i], body, size=body_size)


def fill_cards_by_position(slide, cards, min_abs_left=1000000, body_size=16):
    """按绝对位置配对「添加标题」与正文占位框（支持 GROUP 嵌套）。"""
    title_boxes, body_boxes = collect_card_boxes(slide, min_abs_left)
    for i, (title, body) in enumerate(cards):
        if i < len(title_boxes):
            set_shape_text(title_boxes[i], title, 18, True)
        if i < len(body_boxes):
            fill_shape_body(body_boxes[i], body, size=body_size)
    for sh in title_boxes[len(cards) :]:
        set_shape_text(sh, "", 16)
    for sh in body_boxes[len(cards) :]:
        set_shape_text(sh, "", 16)


def clear_placeholder_titles(slide):
    for sh in list(slide.shapes):
        remove = False
        for sub in walk_shapes([sh]):
            if sub.has_text_frame and sub.text_frame.text.strip() in (
                "点此编辑标题",
                "添加标题",
            ):
                t = sub.text_frame.text.strip()
                if t in ("点此编辑标题", "添加标题"):
                    remove = True
        if remove and sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            sh.element.getparent().remove(sh.element)


def duplicate_slide(pres, index: int) -> int:
    source = pres.slides[index]
    dest = pres.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        el = shp.element
        el.getparent().remove(el)
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")
    return len(pres.slides) - 1


def delete_slide(pres, index: int):
    r_id = pres.slides._sldIdLst[index].rId
    pres.part.drop_rel(r_id)
    del pres.slides._sldIdLst[index]


def move_slide(pres, old_idx: int, new_idx: int):
    slides = pres.slides._sldIdLst
    sld_id = slides[old_idx]
    slides.remove(sld_id)
    slides.insert(new_idx, sld_id)


def prepare_14_slides() -> Presentation:
    shutil.copy(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)
    dup_idx = duplicate_slide(prs, 22)
    assert dup_idx == 27, f"expected dup at 27, got {dup_idx}"
    dup_goals = duplicate_slide(prs, 11)
    assert dup_goals == 28, f"expected dup at 28, got {dup_goals}"

    keep_set = set(KEEP_ORIG)
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in keep_set:
            delete_slide(prs, i)

    remaining = sorted(keep_set)
    for dest in range(len(TARGET_ORDER)):
        want = TARGET_ORDER[dest]
        cur = remaining.index(want)
        if cur != dest:
            move_slide(prs, cur, dest)
            item = remaining.pop(cur)
            remaining.insert(dest, item)
    assert len(prs.slides) == 14
    return prs


# ── 逐页填充 ──────────────────────────────────────────────


def fill_slide_01_cover(slide):
    scrub_shape_text(slide, "LOGO", "")
    replace_exact(slide, "开题报告答辩模板", "广州生产力促进中心中试服务\n管理系统的设计与实现", 36, True)
    replace_exact(slide, "OPENING REPORT DEFENSE TEMPLATE", "Graduation Design Proposal")
    replace_contains(slide, "答辩人", f"答辩人：{STUDENT}    学号：{STUDENT_ID}", 20)
    replace_contains(slide, "指导教师", f"指导教师：{ADVISOR}", 20)
    replace_exact(slide, "202X/XX", DATE, 18)


def fill_slide_02_outline(slide):
    remove_top_level_containing(slide, "1ppt.com")
    scrub_shape_text(slide, "CONTENTES", "CONTENTS")
    scrub_shape_text(slide, "LOGO", "")
    items = [
        "中试背景与政策环境",
        "中心业务现状与系统定位",
        "文献调研与课题切入点",
        "业务活动图：五条主流程",
        "功能模块划分与对应关系",
        "技术路线、进度与成果",
    ]
    nums = ["01", "02", "03", "04"]
    title_shapes = [sh for sh in walk_shapes(slide.shapes) if sh.has_text_frame and sh.text_frame.text.strip() in nums]
    # 找各编号后的标题框
    blocks = []
    for sh in walk_shapes(slide.shapes):
        t = sh.text_frame.text.strip() if sh.has_text_frame else ""
        if t in ("论文选题背景", "论文表达结构", "论文研究方法", "论文主要结论"):
            blocks.append(sh)
    for i, sh in enumerate(blocks[:4]):
        set_shape_text(sh, items[i], size=22, bold=True)
    # 05、06 追加文本框
    if len(blocks) >= 2:
        ref = blocks[-1]
        left, top, w, h = ref.left, ref.top + ref.height + Inches(0.15), ref.width, ref.height
        for j, num in enumerate(["05", "06"], start=4):
            box = slide.shapes.add_textbox(left - Inches(1.2), top + (j - 4) * Inches(0.85), w, h)
            set_shape_text(box, f"{num}  {items[j]}", size=22, bold=True)


def fill_slide_03_zhongshi(slide):
    set_page_title(slide, "什么是中试")
    pairs = [
        ("小试", "实验室阶段验证原理可行性，规模小、投入低"),
        ("中试", "接近生产条件下的放大验证，形成中试报告"),
        ("量产", "规模化生产，设备与资金投入大、试错成本高"),
        ("研究意义", "省略中试环节将导致技术风险向量产阶段传导"),
    ]
    fill_card_pairs(slide, pairs)


def fill_slide_04_policy(slide):
    set_page_title(slide, "政策与产业背景")
    years = ["2024", "2025.05", "2025", "2027"]
    year_shapes = [
        sh
        for sh in walk_shapes(slide.shapes)
        if sh.has_text_frame and sh.text_frame.text.strip() in ("20XX", "2024")
    ]
    year_shapes.sort(key=lambda s: s.left)
    for sh, y in zip(year_shapes[:4], years):
        set_shape_text(sh, y, size=18, bold=True)
    pairs = [
        ("国家与省级布局", "工信部《制造业中试创新发展实施意见》；粤府办〔2024〕7 号"),
        ("广州国家首批", "12 家入选工信部重点培育中试平台（市工信局）"),
        ("平台规模", "全国 2400+ 平台；广州 31 家省首批，数量全省第一"),
        ("建设目标", "2027 年形成多层次制造业中试服务网络"),
    ]
    fill_card_pairs(slide, pairs)
    foot = slide.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(10), Inches(0.4))
    set_shape_text(
        foot,
        "数据来源：工信部、广东省/广州市政府公开信息，2024—2025",
        size=14,
    )


def fill_slide_05_center(slide):
    set_page_title(slide, "中心现状、痛点与本系统")

    # 模板第 16 页：页眉下通栏 + 2×2 宽卡（标题 + 3.26" 正文），适配中心概况与四类痛点。
    overview = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.0), Inches(11.85), Inches(1.05)
    )
    set_section_bullet_list(
        overview,
        [
            (
                "中心概况",
                [
                    "广州生产力促进中心：研发设计、中试服务、技术解决方案",
                    "增城中试基地在运行；累计小试/中试项目 20+",
                    "2025 年高新技术成果产业化试点单位（第三方服务类）",
                ],
            ),
            ("协同角色", ["企业、技术人员、调度员、审核员"]),
        ],
        header_size=15,
        bullet_size=13,
    )

    pains = [
        ("协同方式落后", "电话/微信/Excel 分散协同，缺统一入口与流程状态管理"),
        ("进度不透明", "企业无法实时掌握办理节点，须建立全过程状态追溯"),
        ("复核闭环弱", "报告复核环节易遗漏，反馈模块须双次审核与留痕"),
        ("数据难沉淀", "历史项目数据分散，档案模块承担台账与统计分析"),
    ]
    titles = sorted(
        [
            sh
            for sh in walk_shapes(slide.shapes)
            if sh.has_text_frame and sh.text_frame.text.strip() == "添加标题"
        ],
        key=lambda s: (s.top, s.left),
    )
    bodies = sorted(
        [
            sh
            for sh in walk_shapes(slide.shapes)
            if sh.has_text_frame
            and "单击" in sh.text_frame.text
            and "标题内容" in sh.text_frame.text
        ],
        key=lambda s: (s.top, s.left),
    )
    for i, (title, body) in enumerate(pains):
        if i < len(titles):
            set_shape_text(titles[i], title, 17, True)
        if i < len(bodies):
            set_shape_text(bodies[i], body, 14)
    for sh in titles[len(pains) :]:
        set_shape_text(sh, "", 14)
    for sh in bodies[len(pains) :]:
        set_shape_text(sh, "", 14)

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
    set_shape_text(
        foot,
        "系统范围：中试服务办理流程（申请—归档）；不涵盖实验室试验操作。",
        size=20,
        bold=True,
    )


def fill_panel_groups(slide, panels):
    """按 GROUP 位置填充「点此编辑标题 + 在此输入详细」四宫格（模板第 18 页）。"""
    groups = []
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        title_sh, body_sh = None, None
        for sub, al, at in iter_text_shapes([sh], sh.left, sh.top):
            t = sub.text_frame.text.strip()
            if t == "点此编辑标题":
                title_sh = sub
            elif "在此输入" in t or "详细" in t:
                body_sh = sub
        if title_sh or body_sh:
            groups.append((sh.top, sh.left, title_sh, body_sh))
    groups.sort(key=lambda g: (g[0], g[1]))
    for i, (title, body) in enumerate(panels):
        if i >= len(groups):
            break
        _, _, title_sh, body_sh = groups[i]
        if title_sh:
            set_shape_text(title_sh, title, 22, True)
        if body_sh:
            fill_shape_body(body_sh, body, size=15)
    for j in range(len(panels), len(groups)):
        _, _, title_sh, body_sh = groups[j]
        if title_sh:
            set_shape_text(title_sh, "", 16)
        if body_sh:
            set_shape_text(body_sh, "", 16)


def fill_slide_06_research(slide):
    set_page_title(slide, "研究现状与本课题切入点")
    fill_panel_groups(
        slide,
        [
            (
                "国外研究",
                [
                    "LIMS、项目管理系统较为成熟",
                    "侧重实验数据与内部项目管理",
                    "对政府科技服务机构对外中试及多端接入支持不足",
                ],
            ),
            (
                "国内研究",
                [
                    "成果转化平台、中试基地门户较多",
                    "侧重信息发布、仪器预约与成果展示",
                    "「审核—派单—复核—归档」业务闭环不完整",
                ],
            ),
            (
                "切入点",
                [
                    "明确需求方：广州生产力促进中心",
                    "五条业务主线贯通：需求—评估—调度—反馈—档案",
                    "前后端分离；Web 按角色路由，小程序作移动补充",
                ],
            ),
            (
                "本课题边界",
                [
                    "业务范围：中试服务办理（申请—归档）",
                    "系统边界：流程与数据信息化，不涉及试验执行",
                    "预期成果：可演示系统 + 毕业论文",
                ],
            ),
        ],
    )


def fill_slide_07_goals(slide):
    set_page_title(slide, "研究目标与关键问题")
    for sh, abs_left, _ in iter_text_shapes(slide.shapes):
        if abs_left < 9000000 and sh.has_text_frame and "单击" in sh.text_frame.text:
            set_shape_text(
                sh,
                "设计并实现中试服务全流程信息化管理系统，覆盖需求、评估、调度、"
                "反馈、档案五条业务链路；采用 Web 分角色门户 + 微信小程序移动补充。",
                20,
            )
            break
    for sh in walk_shapes(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip() in ("01", "02", "03", "04"):
            set_shape_text(sh, "", 14)
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.LINE:
            sh.element.getparent().remove(sh.element)
    goal_lbl = slide.shapes.add_textbox(Inches(5.6), Inches(1.25), Inches(2.2), Inches(0.4))
    set_shape_text(goal_lbl, "研究目标", 18, True)
    prob_lbl = slide.shapes.add_textbox(Inches(5.6), Inches(3.45), Inches(2.2), Inches(0.4))
    set_shape_text(prob_lbl, "关键问题", 18, True)
    cards = [
        (
            "目标一",
            [
                "五条业务链线上贯通",
                "节点状态可查询",
                "全过程可追溯",
            ],
        ),
        (
            "目标二",
            [
                "Web 分角色门户",
                "微信小程序移动补充",
                "统一 RESTful 接口",
            ],
        ),
        (
            "问题一",
            [
                "多角色权限划分",
                "责任边界界定",
                "RBAC + 工作流",
            ],
        ),
        (
            "问题二",
            [
                "进度同步机制",
                "复核业务闭环",
                "历史数据入库与统计",
            ],
        ),
    ]
    fill_cards_by_position(slide, cards, min_abs_left=9000000, body_size=14)


def fill_slide_08_activities(slide):
    set_page_title(slide, "五条核心业务活动")
    banner = slide.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(11), Inches(0.45))
    set_shape_text(
        banner,
        "需求管理 → 评估管理 → 调度管理 → 反馈管理 → 档案管理",
        size=20,
        bold=True,
    )
    pairs = [
        (
            "中试需求管理",
            [
                "中试需求提交",
                "需求材料补充",
                "需求受理登记",
                "需求退回通知",
                "图 3-1，10 个活动节点",
            ],
        ),
        (
            "中试评估管理",
            [
                "中试条件评估",
                "技术可行性审",
                "资源需求核定",
                "评估结论出具",
                "图 3-2，12 个活动节点",
            ],
        ),
        (
            "中试调度管理",
            [
                "中试资源匹配",
                "中试任务派发",
                "执行进度跟踪",
                "异常任务重派",
                "图 3-3，11 个活动节点",
            ],
        ),
        (
            "中试反馈管理",
            [
                "试验结果提交",
                "中试报告审核",
                "结果复核确认",
                "中试报告归档",
                "图 3-4，11 个活动节点",
            ],
        ),
        (
            "中试档案管理",
            [
                "项目台账维护",
                "中试周期统计",
                "成功率分析",
                "服务简报生成",
                "图 3-5，9 个活动节点",
            ],
        ),
    ]
    fill_brick_five_cards(slide, pairs, body_size=12)


def fill_slide_09_activity_diagram(slide):
    set_page_title(slide, "图 3-1 中试需求管理业务活动图")
    remove_top_level_containing(slide, "SPORT 02")
    replace_exact(slide, "SPORT 01", "图 3-1 中试需求管理业务活动图", 24, True)
    for sh in list(slide.shapes):
        for sub in walk_shapes([sh]):
            if sub.has_text_frame and "单击此处" in sub.text_frame.text:
                sh.element.getparent().remove(sh.element)
                break
    if os.path.exists(ACTIVITY_IMG):
        slide.shapes.add_picture(ACTIVITY_IMG, Inches(0.6), Inches(1.3), width=Inches(11.8))
    notes = [
        "泳道：上岗位、下活动；三角色：企业、调度员、审核员",
        "主路径：确认—提交—受理—核验—通知—签收—归档—查看",
        "分支：材料不齐退回补充；不同意受理则流程终止",
        "核心活动节点对应模块图「中试需求管理」三级功能",
    ]
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(11.8), Inches(0.9))
    set_shape_text(box, "  |  ".join(notes), size=16)


def fill_slide_10_module_logic(slide):
    set_page_title(slide, "从业务活动到功能模块")
    for sh, abs_left, _ in iter_text_shapes(slide.shapes):
        if abs_left < 9000000 and sh.has_text_frame and "单击" in sh.text_frame.text:
            set_shape_text(
                sh,
                "活动图：描述业务流程与岗位协作关系\n"
                "模块图：按高内聚、低耦合原则划分功能\n"
                "整体结构 1—5—4：1 系统、5 二级模块、各 4 三级功能",
                20,
            )
            break
    cards = [
        ("01 需求", ["提交", "补充", "受理登记"]),
        ("02 评估", ["条件评估", "资源核定", "可行性审", "结论出具"]),
        ("03 调度", ["资源匹配", "任务派发", "进度填报", "异常重派"]),
        ("04 反馈", ["结果提交", "报告审核", "复核确认", "报告归档"]),
    ]
    fill_cards_by_position(slide, cards, min_abs_left=9000000, body_size=14)


def fill_slide_11_module_diagram(slide):
    set_page_title(slide, "功能模块图（1—5—4 结构）")
    remove_top_level_containing(slide, "SPORT 02")
    replace_exact(slide, "SPORT 01", "功能模块图（1—5—4 结构）", 24, True)
    for sh in list(slide.shapes):
        for sub in walk_shapes([sh]):
            if sub.has_text_frame and "单击此处" in sub.text_frame.text:
                sh.element.getparent().remove(sh.element)
                break
    if os.path.exists(MODULE_IMG):
        slide.shapes.add_picture(MODULE_IMG, Inches(0.3), Inches(1.2), width=Inches(12.5))
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.4))
    set_shape_text(cap, "五个二级模块与五张业务活动图一一对应，模块间保持低耦合", size=18)


def add_round_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    return shape


def _shape_bottom(sh):
    return sh.top + sh.height


def _shape_right(sh):
    return sh.left + sh.width


def _rects_overlap(a, b, margin=0):
    return not (
        _shape_right(a) + margin <= b.left
        or _shape_right(b) + margin <= a.left
        or _shape_bottom(a) + margin <= b.top
        or _shape_bottom(b) + margin <= a.top
    )


def verify_slide_12_layout(slide) -> list[str]:
    """检查第 12 页内容区是否越界或互相重叠。"""
    issues = []
    content = [
        sh
        for sh in slide.shapes
        if sh.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE)
        and sh.top >= Inches(1.0)
    ]
    for sh in content:
        if _shape_right(sh) > SLIDE_W + Inches(0.05):
            issues.append(f"形状越界右缘: top={sh.top / 914400:.2f}in")
        if _shape_bottom(sh) > SLIDE_H + Inches(0.05):
            issues.append(f"形状越界底缘: top={sh.top / 914400:.2f}in")
    for i, a in enumerate(content):
        for b in content[i + 1 :]:
            if _rects_overlap(a, b):
                issues.append(
                    f"内容重叠: T={a.top / 914400:.2f} 与 T={b.top / 914400:.2f}"
                )
    return issues


def set_title_body_box(shape, title: str, body: str, title_size=20, body_size=18):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(title_size)
    r1.font.name = "微软雅黑"
    r1.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(body_size)
    r2.font.name = "微软雅黑"


def fill_slide_12_architecture(slide):
    set_page_title(slide, "系统架构与技术选型")
    clear_placeholder_titles(slide)
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.FREEFORM and sh.has_text_frame:
            sh.element.getparent().remove(sh.element)

    layers = [
        "表现层 · Web（Vue3 + Element Plus，按登录角色路由）",
        "表现层 · 微信小程序（企业/技术人员移动补充）",
        "接口层 · RESTful API（前后端分离）",
        "业务层 · Spring Boot + MyBatis",
        "数据层 · MySQL 8.0",
    ]
    arch_top = Inches(1.15)
    layer_h = Inches(0.30)
    layer_gap = Inches(0.08)
    layer_w = Inches(9.2)
    layer_left = Inches(2.07)

    for i, text in enumerate(layers):
        top = arch_top + i * (layer_h + layer_gap)
        bar = add_round_rect(slide, layer_left, top, layer_w, layer_h, ARCH_FILL, ARCH_LINE)
        set_shape_text(bar, text, size=14, bold=True)
        bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cols = [
        (
            "Web · 分角色门户",
            [
                "统一 Vue 前端，RBAC 角色路由",
                "企业端：需求提交、进度查询、反馈接收",
                "中心端：审核、派单、复核",
                "技术人员：任务接单与填报",
            ],
        ),
        (
            "微信小程序",
            [
                "外出及现场场景移动补充",
                "进度查询、消息通知",
                "精简填报",
                "复杂操作集中于 Web 端",
            ],
        ),
        (
            "后端与数据",
            [
                "Spring Boot 业务逻辑层",
                "MyBatis 持久化 + MySQL",
                "按领域划分模块",
                "接口层与表现层解耦",
            ],
        ),
    ]
    col_w = Inches(3.75)
    col_h = Inches(3.20)
    col_gap = Inches(0.54)
    col_top = Inches(3.20)
    col_start = Inches(0.50)

    for i, (title, body) in enumerate(cols):
        left = col_start + i * (col_w + col_gap)
        card = add_round_rect(slide, left, col_top, col_w, col_h, CARD_FILL, CARD_LINE)
        set_title_bullet_box(card, title, body, title_size=18, bullet_size=14)


def fill_slide_13_schedule(slide):
    set_page_title(slide, "进度安排与已有工作")
    cards = [
        ("1—2 调研与开题", "需求调研、开题报告撰写"),
        ("3—4 业务建模", "五张活动图 + 模块图 ✓ 已完成"),
        ("5—7 后端开发", "库表设计、核心业务 API"),
        ("8—10 Web 端", "中心管理端 + 企业/技术人员门户"),
        ("11—12 小程序", "移动侧查进度、填报、消息"),
        ("13—16 测试与论文", "联调测试、可演示系统、论文定稿"),
    ]
    fill_cards_by_position(slide, cards, min_abs_left=0)
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5))
    set_shape_text(
        foot,
        "当前进度：业务活动图 5 张、功能模块图 1 张已完成；可运行系统于开题后进入开发阶段。",
        size=18,
        bold=True,
    )


def fill_slide_14_closing(slide):
    replace_exact(slide, "LOGO", "")
    scrub_shape_text(slide, "LOGO", "")
    replace_exact(slide, "谢谢观看", "恳请各位老师批评指正", 36, True)
    replace_exact(slide, "OPENING REPORT DEFENSE TEMPLATE", "")
    replace_exact(slide, "202X/XX", DATE, 18)
    replace_contains(slide, "答辩人", f"答辩人：{STUDENT}    学号：{STUDENT_ID}", 20)
    replace_contains(slide, "指导教师", f"指导教师：{ADVISOR}", 20)
    summary = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(2))
    set_shape_text(
        summary,
        "选题依据：中试政策驱动 + 中心真实业务需求，系统边界明确\n"
        "方案阶段：五张活动图与功能模块图已对齐，可指导后续开发\n"
        "实施计划：16 周内完成 Web + 小程序系统及毕业论文\n"
        "预期成果：可演示 Web 分角色门户 + 小程序移动补充 + 毕业论文",
        size=22,
    )


FILLERS = [
    fill_slide_01_cover,
    fill_slide_02_outline,
    fill_slide_03_zhongshi,
    fill_slide_04_policy,
    fill_slide_05_center,
    fill_slide_06_research,
    fill_slide_07_goals,
    fill_slide_08_activities,
    fill_slide_09_activity_diagram,
    fill_slide_10_module_logic,
    fill_slide_11_module_diagram,
    fill_slide_12_architecture,
    fill_slide_13_schedule,
    fill_slide_14_closing,
]

EXPECTED_TITLES = [
    "广州生产力促进中心",
    "01",
    "什么是中试",
    "政策与产业背景",
    "中心现状",
    "研究现状",
    "研究目标",
    "五条核心",
    "图 3-1",
    "从业务活动",
    "功能模块图",
    "系统架构",
    "进度安排",
    "恳请各位",
]


def verify(prs: Presentation) -> list[str]:
    issues = []
    if len(prs.slides) != 14:
        issues.append(f"页数应为 14，实际 {len(prs.slides)}")

    all_text = []
    for i, slide in enumerate(prs.slides, 1):
        page_text = " ".join(
            sh.text_frame.text
            for sh in walk_shapes(slide.shapes)
            if sh.has_text_frame and sh.text_frame.text.strip()
        )
        all_text.append(page_text)
        if "1ppt.com" in page_text.lower():
            issues.append(f"第{i}页仍含 1ppt.com 广告链接")
        for m in PLACEHOLDER_MARKERS:
            if m in page_text:
                issues.append(f"第{i}页可能残留占位符「{m}」")
                break

    if STUDENT not in all_text[0]:
        issues.append("封面缺少答辩人姓名")
    if ADVISOR not in all_text[0]:
        issues.append("封面缺少指导教师")
    if "2400+" not in all_text[3] and "2400" not in all_text[3]:
        issues.append("第4页缺少政策数据")
    if "数据来源" not in all_text[3]:
        issues.append("第4页缺少数据来源脚注")
    if "SPORT 02" in " ".join(all_text):
        issues.append("图页仍含 SPORT 02 占位")
    if "业务活动图" not in all_text[8] and "需求管理" not in all_text[8]:
        issues.append("第9页缺少活动图标题")
    if "功能模块" not in all_text[10]:
        issues.append("第11页缺少模块图标题")
    if "✓" not in all_text[12] and "已完成" not in all_text[12]:
        issues.append("第13页未标注已完成工作")

    joined = " ".join(all_text)
    if "双端协同" in joined or "可演示双端系统" in joined:
        issues.append("仍含过时表述「双端」")
    if "Web 分角色" not in joined and "分角色门户" not in joined:
        issues.append("缺少 Web 分角色门户表述")
    if "移动补充" not in joined:
        issues.append("缺少小程序移动补充表述")

    if len(prs.slides) >= 12:
        for msg in verify_slide_12_layout(prs.slides[11]):
            issues.append(f"第12页排版: {msg}")

    if os.path.exists(ACTIVITY_IMG) and os.path.exists(MODULE_IMG):
        pic_counts = []
        for slide in prs.slides:
            pic_counts.append(
                sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
            )
        if pic_counts[8] < 1:
            issues.append("第9页未插入活动图图片")
        if pic_counts[10] < 1:
            issues.append("第11页未插入模块图图片")

    return issues


def main():
    prs = prepare_14_slides()
    for slide, filler in zip(prs.slides, FILLERS):
        filler(slide)
        scrub_placeholders(slide)
    prs.save(OUTPUT)

    prs2 = Presentation(OUTPUT)
    issues = verify(prs2)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs2.slides)}")
    print("--- 总览检查 ---")
    if issues:
        for x in issues:
            print(f"  [WARN] {x}")
        print("结论: 有警告，请人工复核")
    else:
        print("  全部检查项通过")
    return 0 if not issues else 1


if __name__ == "__main__":
    raise SystemExit(main())
