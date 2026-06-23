# -*- coding: utf-8 -*-
"""Generate opening defense PPT for topic 2."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

DIR = os.path.dirname(os.path.dirname(__file__))
OUT = os.path.join(DIR, "答辩PPT", "开题答辩-广州生产力促进中心中试服务管理系统.pptx")
IMG_DIR = os.path.join(DIR, "业务活动图")
MODULE_IMG = os.path.join(DIR, "功能模块图", "广州生产力促进中心中试服务管理系统的设计与实现.png")

TITLE_COLOR = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)


def set_title(slide, text, subtitle=None):
    slide.shapes.title.text = text
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_bullets(text_frame, items, size=22):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "微软雅黑"


def add_content_slide(prs, title, bullets, size=22):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    add_bullets(tf, bullets, size)
    return slide


def main():
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "毕业设计开题答辩"
    slide.placeholders[1].text = (
        "题目：广州生产力促进中心中试服务管理系统的设计与实现\n"
        "专业：软件工程\n"
        "答辩人：________\n"
        "指导教师：________\n"
        "2026年6月"
    )

    slides_data = [
        ("汇报提纲", [
            "选题背景与意义",
            "研究目标与内容",
            "业务痛点分析",
            "三个核心业务活动",
            "系统架构与功能模块",
            "技术路线与进度安排",
        ]),
        ("选题背景", [
            "科技成果转化需经历小试、中试、量产三个阶段",
            "中试（中间试验）是验证技术能否放大试产的关键环节",
            "广州生产力促进中心承担中试服务等科技服务职能",
            "中心现有协同依赖电话、微信群、Excel，缺乏系统闭环",
        ]),
        ("研究意义", [
            "理论：将软件工程方法应用于科技服务机构业务流程数字化",
            "实践：提升中试需求受理、派单执行、报告复核的效率与透明度",
            "社会：支撑广州科技成果转移转化与中试服务平台建设",
        ]),
        ("业务痛点", [
            "协同方式落后：电话、微信群、Excel 难以长期稳定运营",
            "状态不透明：企业无法实时掌握中试办理进度",
            "责任不清晰：派单、执行、复核环节容易脱节",
            "闭环不完整：报告复核与历史归档容易缺失",
            "数据无沉淀：难以支撑服务效率分析与决策",
        ]),
        ("系统定位（做什么）", [
            "为广州生产力促进中心建设「中试服务全流程管理系统」",
            "管理申请 → 审核 → 评估 → 派单 → 执行 → 反馈 → 复核 → 归档",
            "系统管流程与协同，不替代实验室内的中试试验操作本身",
            "多端架构：Web 分角色门户 + 微信小程序移动补充",
        ]),
        ("四个系统角色", [
            "企业项目负责人（Web 门户为主 + 小程序补充）：提交需求、补充材料、查看进度",
            "中试技术人员（Web 门户为主 + 小程序补充）：接收任务、执行反馈、提交结果",
            "中试调度员（Web）：受理审核、条件评估、资源匹配、派单督办",
            "中试审核员（Web）：报告复核、结果确认、项目归档",
        ]),
        ("核心活动1：需求提交与受理审核", [
            "企业提交中试需求 → 调度员受理登记",
            "判断：材料是否齐全合规",
            "否 → 退回通知 → 企业补充材料 → 重新提交（循环）",
            "是 → 受理通过，进入评估环节",
            "对应模块：中试需求管理",
        ]),
        ("核心活动2：条件评估与任务调度", [
            "调度员开展中试条件评估",
            "判断：是否具备中试条件",
            "否 → 出具评估结论，需求终止",
            "是 → 匹配资源、派发任务 → 技术人员执行",
            "判断：是否异常需重派 → 是则重派（循环）",
            "对应模块：中试评估管理、中试调度管理",
        ]),
        ("核心活动3：结果反馈与复核归档", [
            "技术人员提交试验结果与中试报告",
            "审核员审核报告材料",
            "判断：报告是否复核合格",
            "否 → 退回修改 → 重新提交（循环）",
            "是 → 报告归档 → 更新项目台账",
            "对应模块：中试反馈管理、中试档案管理",
        ]),
        ("系统总体架构", [
            "表现层：Web 端（Vue + Element Plus，按账号分角色页面）+ 微信小程序（移动补充）",
            "业务层：Spring Boot 后端服务",
            "数据层：MySQL 关系型数据库",
            "多端通过 RESTful API 协同，形成完整业务闭环",
        ]),
        ("功能模块设计（5大模块）", [
            "中试需求管理：需求提交 / 材料补充 / 受理登记 / 退回通知",
            "中试评估管理：条件评估 / 可行性审 / 资源核定 / 结论出具",
            "中试调度管理：资源匹配 / 任务派发 / 进度跟踪 / 异常重派",
            "中试反馈管理：结果提交 / 报告审核 / 复核确认 / 报告归档",
            "中试档案管理：台账维护 / 周期统计 / 成功率分析 / 简报生成",
        ]),
        ("技术路线", [
            "第1—2周：需求调研、开题答辩",
            "第3—4周：业务建模（业务活动图、功能模块图）",
            "第5—7周：数据库与后端开发",
            "第8—10周：Web 端（中心管理 + 企业/技术人员门户）",
            "第11—12周：微信小程序（移动补充）",
            "第13—16周：联调测试、论文撰写与答辩",
        ]),
        ("预期成果", [
            "可运行的 Web 分角色门户 + 微信小程序系统",
            "毕业设计论文 1 篇",
            "业务活动图、功能模块图、用例图、E-R 图等配套文档",
            "系统演示与答辩材料",
        ]),
        ("谢谢！", [
            "敬请各位老师批评指正",
            "",
            "Q & A",
        ]),
    ]

    for title, bullets in slides_data:
        add_content_slide(prs, title, bullets)

    # 插入业务活动图（3页）
    activity_imgs = [
        ("业务活动图1-中试需求提交与受理审核.png", "业务活动图：需求提交与受理审核"),
        ("业务活动图2-中试条件评估与任务调度.png", "业务活动图：条件评估与任务调度"),
        ("业务活动图3-试验结果反馈与报告复核归档.png", "业务活动图：结果反馈与复核归档"),
    ]
    for fname, title in activity_imgs:
        path = os.path.join(IMG_DIR, fname)
        if os.path.exists(path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            slide.shapes.add_picture(path, Inches(0.4), Inches(1.2), width=Inches(12.5))

    # 插入功能模块图
    if os.path.exists(MODULE_IMG):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "功能模块图"
        slide.shapes.add_picture(MODULE_IMG, Inches(0.3), Inches(1.2), width=Inches(12.7))

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
